from tkinter import *
from tkinter import filedialog
import tkinter.messagebox as msgbox
import tkinter.ttk as ttk
import os
import datetime
import subprocess
import pptx


"""프로그래밍 순서
0. gui의 형태를 먼저 스케치해야함
1. root와 frame, 기초적인 버튼등을 배치 (레이아웃)
2. 기본적인 함수를 정의 (버튼 등이 동작하도록)
3. 기능적인 함수를 정의(기능을 요구하는 함수) 3.1, 3.2 등의 저장을 요구

"""
# sys.stdout = open('log.txt', 'a')


RootName = "AutoSlide"
root = Tk()
root.title(RootName)

# 추가 필요 : 더블클릭으로 리스트에서 다음 리스트로 추가할 수 있도록 설정할 수 있으면 좋겠다.



#기본경로설정
if os.path.isfile("db1.txt"): #db1file에서 최근경로를 가져옴
    f= open("db1.txt", 'r')
    folder_selected = f.readline()
    f.close()
else: #db1file이 존재하지 않는경우, 새로운 파일 열기
    f= open("db1.txt", 'w')
    f.close()
    folder_selected = ""

#기본변수 설정
pptmain_selected = ""



#definition
def dir_pathfind(): #dir 찾아보기버튼
    global folder_selected
    folder_selected = filedialog.askdirectory()
    if folder_selected == '':
        return
    txt_dir_path.delete(0,END)
    txt_dir_path.insert(END, folder_selected)
def lst_reload(): #dir 확인버튼, 최적화x
    folder_selected = txt_dir_path.get()
    try:
        global pptlist
        pptlist = []
        lstb_l.delete(0,END)
        li = [k for k in os.listdir(folder_selected) if os.path.isfile(os.path.join(folder_selected, k))]
        for i in li:
            if i[-5:] == ".pptx":
                pptlist.append(i[:-5])
            else:
                pass
        for ppt in pptlist:
            lstb_l.insert(END, ppt)
        with open("db1.txt", 'w') as f:
            f.write(folder_selected)
    except Exception as err:
        msgbox.showerror("에러", err)
def lst_reload1(): #dir 확인버튼, 최적화o
    try:
        folder_selected = filedialog.askdirectory()
        if not folder_selected:
            return
        pptlist = [os.path.splitext(file)[0] for file in os.listdir(folder_selected) if file.endswith(".pptx")]
        lstb_l.delete(0, END)
        for ppt in pptlist:
            lstb_l.insert(END, ppt)
        with open("db1.txt", 'w') as f:
            f.write(folder_selected)
    except Exception as err:
        msgbox.showerror(RootName, err)

#좌측버튼
def slide_addppt(): #lstb_r로 항목이동
    num_listppt = lstb_l.curselection()
    for i in num_listppt:
        a =lstb_l.get(i)
        lstb_r.insert(END, a)
def slide_lopenfile(): #lstb_l의 선택된 파일 열기
    num_listppt = lstb_l.curselection()
    answer = msgbox.askyesno(RootName, "해당 파일을 외부 프로그램으로 실행합니다.\n여러 프로그램을 한번에 실행 할 경우 성능저하가 발생합니다.")
    if answer:
        for i in num_listppt:
                a =lstb_l.get(i)
                file_path = folder_selected+"/"+a+".pptx"
                subprocess.run(['start', '', file_path], shell=True)
    else: pass

#우측버튼
def fileselect1(): #fileselect, 교안파일열기, 상위
    global pptmain_selected
    pptmain_selected = filedialog.askopenfilename(title="파일을 선택하세요", 
                filetypes=(("powerpoint 파일", "*.pptx"), ("모든 파일", "*.*")),
                initialdir= "Documents/카카오톡 받은 파일")
    if pptmain_selected[-5:] == ".pptx": 
        a = pptmain_selected[::-1]
        say =  f"설교 파일 [{pptmain_selected[-a.find("/"):]}]"
        lstb_r.insert(0, say)
    else:
        msgbox.showerror(RootName, "(*.pptx) powerpoint 파일이 선택되지 않았습니다.\n다시 시도하십시오.")
def slide_searchppt1(): #fileselect, 교안파일열기, 파일존재여부 확인
    if pptmain_selected == "":
        fileselect1()        
    else:
        answer = msgbox.askretrycancel(
            RootName, 
            "이미 선택된 파일이 존재합니다. \n다른 파일로 다시시도합니까?")
        if answer:
            lstb_r.delete(0)
            fileselect1()
def slide_delppt1(): #교안파일제거
    global pptmain_selected
    if pptmain_selected == "":
        pass
    else:
        lstb_r.delete(0)
        pptmain_selected = ""
def slide_searchppt2(): #fileselect, 가사파일 열기 ... 구현할 수 없을 듯하여 폐기
    file = filedialog.askopenfilename(title="파일을 선택하세요", 
                filetypes=(("powerpoint 파일", "*.pptx"), ("모든 파일", "*.*")),
                initialdir= "Documents/카카오톡 받은 파일")
    if file.endswith('.pptx'): 
        a = file[::-1]
        lstb_r.insert(END, file[-a.find("/"):])
    else:
        msgbox.showerror(RootName, "(*.pptx) powerpoint 파일이 선택되지 않았습니다.\n다시 시도하십시오.")
def slide_ropenfile(): #lstb_r의 선택된 파일 열기
    num_listppt = lstb_r.curselection()
    answer = msgbox.askyesno(RootName,"해당 파일을 외부 프로그램으로 실행합니다.\n여러 프로그램을 한번에 실행 할 경우 성능저하가 발생합니다.")
    if answer:
        for i in num_listppt:
            a =lstb_r.get(i)
            file_path = folder_selected+"/"+a+".pptx"
            subprocess.run(['start', '', file_path], shell=True)
    else: pass
def slide_delppt2(): #lstb_r의 선택된 파일 제거
    for index in reversed(lstb_r.curselection()):
        if index ==0 and pptmain_selected != "":
            pass
        else:
            lstb_r.delete(index)
    
def lstb_raise():
    selected_indices = lstb_r.curselection()
    # print("Index to raise:", index_to_raise)
    # lstb_r.lift(index_to_raise)
    if selected_indices:
        # Get the index of the first selected item
        index_to_raise = selected_indices[0]
        print("Index to raise:", index_to_raise)
        # Raise the item in the Listbox
        lstb_r.lift(index_to_raise)

def lstb_lower():
    selected_indices = lstb_r.curselection()

    if selected_indices:
        # Get the index of the last selected item
        index_to_lower = selected_indices[-1]

        # Lower the item in the Listbox
        lstb_r.lower(index_to_lower)

def slide_start():
    pass


#슬라이드 합치기
def lstb_getitem(lstb): #listbox의 모든 item을 list형태로 가져옴
    allitems= [lstb.get(index) for index in range(lstb.size())] #range 함수!!!
    return allitems

def slide_start():
    try:
        lst_mainfile = lstb_getitem(lstb_r)
        for file in lst_mainfile:
            file_path = folder_selected+"/"+file+".pptx"
        
        







        #fime name
        # 현재 날짜를 얻어옴
        current_date = datetime.date.today()
        # 날짜 형식 지정 (예: "년월일.txt")
        date_format = "{:%Y%m%d}.txt"
        formatted_file_name = date_format.format(current_date)

            
            




    except Exception as err:
        msgbox.showerror("에러", err)
    


#

def clone_slide(prs, slide):
    new_slide = prs.slides.add_slide(slide.slide_layout)
    for shape in slide.shapes:
        new_shape = new_slide.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)
        new_shape.text = shape.text

    return new_slide






















#dir frame
dir_frame = LabelFrame(root, text="Directory Select")
dir_frame.pack(fill="x", padx=5, pady=5)

txt_dir_path = Entry(dir_frame)
txt_dir_path.pack(side="left", fill="x", expand=True, padx=5, pady=5, ipady=4)
txt_dir_path.insert(END, folder_selected)

btn_lst_reload = Button(dir_frame, text="확인", width=10, command=lst_reload)
btn_lst_reload.pack(side="right", padx=5, pady=5)


btn_dir_path = Button(dir_frame, text="찾아보기", width=10, command=dir_pathfind)
btn_dir_path.pack(side="right", padx=5, pady=5)


#combineframe1
combineframe1 = Frame(root)
combineframe1.pack(fill="x", padx=5, pady=5)

#list frame
lst_frame1 = LabelFrame(combineframe1, text="File List")
lst_frame1.pack(side="left", padx=5, pady=5)

lst_frame2 = Frame(lst_frame1)
lst_frame2.pack(padx=5, pady=5)

lst_scroll = Scrollbar(lst_frame2)
lst_scroll.pack(side="right", fill="y")

lstb_l = Listbox(lst_frame2, selectmode="extended", height=25, yscrollcommand=lst_scroll.set, width=43)
lstb_l.pack(side="left", fill="both", expand="True")
lst_scroll.config(command=lstb_l.yview)

btn_lst_select = Button(lst_frame1, text="선택 적용", width=10, command= slide_addppt)
btn_lst_select.pack(side="right", padx=5, pady=5)

btn_lst_select = Button(lst_frame1, text="선택 실행", width=10, command= slide_lopenfile)
btn_lst_select.pack(side="right", padx=5, pady=5)

#selected ppt frame
slc_frame = LabelFrame(combineframe1, text="Selected Files")
slc_frame.pack(side="right", padx=5, pady=5, fill="y")

lstb_r = Listbox(slc_frame, height=5, selectmode="extended", width=43)
lstb_r.pack(fill="both", expand="True", padx=5, pady=5, side="left")

#btn ppt.
combineframe2 = Frame(slc_frame)
combineframe2.pack(side="right")
btn_slc_search = Button(combineframe2, text="설교 파일", width=10, command = slide_searchppt1)
btn_slc_search.pack(padx=5, pady=5)
btn_slc_search = Button(combineframe2, text="설교 제거", width=10, command = slide_delppt1)
btn_slc_search.pack(padx=5, pady=5)
# lbl_pptframe = Label(combineframe2, text="")
# lbl_pptframe.pack(padx=5, pady=5)
# btn_slc_search = Button(combineframe2, text="파일 추가", width=10, command = slide_searchppt2)
# btn_slc_search.pack(padx=5, pady=5)
btn_slc_del = Button(combineframe2, text="선택 열기", width=10, command = slide_ropenfile)
btn_slc_del.pack(padx=5, pady=5)
btn_slc_del = Button(combineframe2, text="선택 삭제", width=10, command = slide_delppt2)
btn_slc_del.pack(padx=5, pady=5)
lbl_pptframe = Label(combineframe2, text="")
lbl_pptframe.pack(padx=5, pady=5)
btn_slc_goup = Button(combineframe2, text="위로 이동", width=10, command = lstb_raise)
btn_slc_goup.pack(padx=5, pady=5)
btn_slc_godown = Button(combineframe2, text="아래로 이동", width=10, command = lstb_lower)
btn_slc_godown.pack(padx=5, pady=5)










#option frame
option_frame = LabelFrame(root, text="Option")
option_frame.pack(padx=5, pady=5, fill="x")

#1. 중간 슬라이드 개수 옵션
lbl_insslide = Label(option_frame, text="삽입슬라이드", width=10)
lbl_insslide.pack(padx=7, pady=7, side="left")

option_insslide = ["1개", "2개", "3개"]
cmb_insslide = ttk.Combobox(option_frame, state="readonly", values=option_insslide, width=10)
cmb_insslide.current(0)
cmb_insslide.pack(padx=7, pady=7, side="left")

#2. 우상단 워터마크 옵션
lbl_marking = Label(option_frame, text="워터마크", width=9)
lbl_marking.pack(side="left", padx=7, pady=7)

option_useornot = ["사용", "제거"]
cmb_marking = ttk.Combobox(option_frame, state="readonly", values=option_useornot, width=10)
cmb_marking.current(1)
cmb_marking.pack(side="left", padx=7, pady=7)

#3. 애니메이션 옵션
lbl_animation = Label(option_frame, text="에니매이션", width=9)
lbl_animation.pack(side="left", padx=7, pady=7)

cmb_animation = ttk.Combobox(option_frame, state="readonly", values=option_useornot, width=10)
cmb_animation.current(1)
cmb_animation.pack(side="left", padx=7, pady=7)

#4. 테마옵션
option_theme = ["기본테마", "테마1", "테마2"]
lbl_animation = Label(option_frame, text="테마선택", width=9)
lbl_animation.pack(side="left", padx=7, pady=7)

cmb_animation = ttk.Combobox(option_frame, state="readonly", values=option_theme, width=10)
cmb_animation.current(0)
cmb_animation.pack(side="left", padx=7, pady=7)


#progressbar frame
frame_progress = LabelFrame(root, text="Progress")
frame_progress.pack(fill="x", padx=5, pady=5, ipady=5)

p_var = DoubleVar()
progress_bar = ttk.Progressbar(frame_progress, maximum=100, variable=p_var)
progress_bar.pack(fill="x", padx=5, pady=5)


#run frame
frame_run = Frame(root)
frame_run.pack(fill="x", padx=5, pady=5, ipady=5)

btn_close = Button(frame_run, padx=5, pady=5, text="닫기", width=12, command = root.quit)
btn_close.pack(side="right", padx=5, pady=5)

btn_start = Button(frame_run, padx=5, pady=5, text="시작", width=12, command = slide_start)
btn_start.pack(side="right", padx=5, pady=5)










#이후 동작관련

# #ppt목록 자동으로 불러오기 (작동과 동시에 디렉토리가 없으면 종료되는 문제가 있음)
# if folder_selected != "":
#     folder_selected = txt_dir_path.get()
#     try:
#         if folder_selected :
#             pass
#         else:
#             folder_selected = ""
#         global pptlist
#         pptlist = []
#         lstb_l.delete(0,END)
#         li = [k for k in os.listdir(folder_selected) if os.path.isfile(os.path.join(folder_selected, k))]
#         for i in li:
#             if i[-5:] == ".pptx":
#                 pptlist.append(i[:-5])
#             else:
#                 pass
#         for ppt in pptlist:
#             lstb_l.insert(END, ppt)
#         with open("db1.txt", 'w') as f:
#             f.write(folder_selected)
#     except Exception as err:
#         msgbox.showerror("에러", err)





root.resizable(False, False)
root.mainloop( )

# sys.stdout.close()