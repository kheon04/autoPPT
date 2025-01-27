from pptx import Presentation


# slide number setting
SLD_LAYOUT_TITLE_AND_CONTENT = 1

# Default Adding Slide Methods
# By appending add_slide()
prs = Presentation()
slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
slide = prs.slides.add_slide(slide_layout)

# Insert 기능은 존재하지 않음.




##2. Shapes
# AutoShapes(shape shapes, text boxes, placeholders)
# line/connector
# picture
# graphic frame(table, chart, smart art, media clip)

# slide는, prs에 추가했던 슬라이드... 이에 존재하는 모든 shapes의 list를 의미
shapes = slide.shapes




##2_1. AutoShapes
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Default Adding AutoShapes
left = top = width = height = Inches(1.0)
shape = shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
)

# Inches, cm, pt
print(Inches(1).cm, Inches(1).pt)
length = Inches(1)
print(length.cm, length.pt, length)
#(Inches(1) == 914400, EMU로 저장된 값이므로, 이는 어떤 이유가 있으나 알 필요X)

# position, demention()
print(shape.left, shape.left.inches)

# Fill (docs 참고)
# Line (outLine, docs 참고)

# Adjusting an autoshape
callout_sp = shapes.add_shape(
    MSO_SHAPE.LINE_CALLOUT_2_ACCENT_BAR, left, top, width, height
)
# get the callout line coming out of the right place
adjs = callout_sp.adjustments
adjs[0] = 0.5   # vert pos of junction in margin line, 0 is top
adjs[1] = 0.0   # horz pos of margin ln wrt shape width, 0 is left side
adjs[2] = 0.5   # vert pos of elbow wrt margin line, 0 is top
adjs[3] = -0.1  # horz pos of elbow wrt shape width, 0 is margin line
adjs[4] = 3.0   # vert pos of line end wrt shape height, 0 is top
a5 = adjs[3] - (adjs[4] - adjs[0]) * height/width
adjs[5] = a5    # horz pos of elbow wrt shape width, 0 is margin line

# rotate 45 degrees counter-clockwise
callout_sp.rotation = -45.0





##3. Placeholders
# access by idx (0~5(기본), 10~(사용자 지정))
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[8])
for shape in slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))

# Identify and Characterize
# shape type이 어떤 도형이건 간에, MSO_SHAPE_TYPE.PLACEHOLDER로 고정됨
# 또한 모든 shape도, placeholder_format을 사용할 수 있으나 ---> 만약 placeholder가 아니라면, valueerror를 발생시킴
# if문을 활용, placeholder인지를 확인하는 과정을 거쳐야한다
for shape in slide.shapes:
    if shape.is_placeholder:
        phf = shape.placeholder_format
        print('%d, %s' % (phf.idx, phf.type))

# Insert content
# Picture
slide = prs.slides.add_slide(prs.slide_layouts[8])
placeholder = slide.placeholders[1]
print(placeholder.name)
print(placeholder.placeholder_format.type)
# Table (docs 참조)
# Chart (docs 참조)
# Title, (Subtitle은 불가능)
title_placeholder = slide.shapes.title
title_placeholder.text = 'Air-speed Velocity of Unladen Swallows'





##4. Working with Text
# Accessing
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text_frame = shape.text_frame

# Adding text (간략화된 방식)
shape.text = 'foobar'

# is equivalent to ... 
text_frame = shape.text_frame
text_frame.clear()
p = text_frame.paragraphs[0]
run = p.add_run()
run.text = 'foobar'

# 원초적인 방식. (paragraph와, text 추가)
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'
tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'
p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p = tf.add_paragraph()  # 해당 경우, 단락이 추가되지만,
p.text = 'test1\ntest2' # \n을 이용한 enter는 새로운 단락을 추가하지 않음 (shift+enter와 같은 역할)

# for문을 이용한 간략화 방식 (paragraph와, text 추가)
paragraph_strs = [
    'Egg, bacon, sausage and spam.',
    'Spam, bacon, sausage and spam.',
    'Spam, egg, spam, spam, bacon and spam.'
]

text_frame = shape.text_frame
text_frame.clear()  # remove any existing paragraphs, leaving one empty one

p = text_frame.paragraphs[0]
p.text = paragraph_strs[0]

for para_str in paragraph_strs[1:]:
    p = text_frame.add_paragraph()
    p.text = para_str

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
# font
font = run.font
font.name = 'Calibri'
font.size = Pt(18)
font.bold = True
font.italic = None  # cause value to be inherited from theme
font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

# 기타.
font.color.rgb = RGBColor(0xFF, 0x7F, 0x50) # 색 지정
run.hyperlink.address = 'https://github.com/scanny/python-pptx' # 하이퍼링크 지정


##5. Slides
# Slide masters
slide_master = prs.slide_masters[0]

# Slide layouts
title_slide_layout = prs.slide_masters[0].slide_layouts[0]

# Slides
first_slide = prs.slides[0]
