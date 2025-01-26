from pptx import Presentation


# slide number setting
SLD_LAYOUT_TITLE_AND_CONTENT = 1

# Default Adding Slide Methods
# By appending add_slide()
prs = Presentation()
slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
slide = prs.slides.add_slide(slide_layout)

# Insert 기능은 존재하지 않음.

##2. Shapes
# AutoShapes(shape shapes, text boxes, placeholders)
# line/connector
# picture
# graphic frame(table, chart, smart art, media clip)

# slide는, prs에 추가했던 슬라이드... 이에 존재하는 모든 shapes의 list를 의미
shapes = slide.shapes

##2_1. AutoShapes
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Default Adding AutoShapes
left = top = width = height = Inches(1.0)
shape = shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
)

# Inches, cm, pt
print(Inches(1).cm, Inches(1).pt)
length = Inches(1)
print(length.cm, length.pt, length)
#(Inches(1) == 914400, EMU로 저장된 값이므로, 이는 어떤 이유가 있으나 알 필요X)

# position, demention()
print(shape.left, shape.left.inches)

# Fill (docs 참고)
# Line (outLine, docs 참고)

# Adjusting an autoshape
callout_sp = shapes.add_shape(
    MSO_SHAPE.LINE_CALLOUT_2_ACCENT_BAR, left, top, width, height
)
# get the callout line coming out of the right place
adjs = callout_sp.adjustments
adjs[0] = 0.5   # vert pos of junction in margin line, 0 is top
adjs[1] = 0.0   # horz pos of margin ln wrt shape width, 0 is left side
adjs[2] = 0.5   # vert pos of elbow wrt margin line, 0 is top
adjs[3] = -0.1  # horz pos of elbow wrt shape width, 0 is margin line
adjs[4] = 3.0   # vert pos of line end wrt shape height, 0 is top
a5 = adjs[3] - (adjs[4] - adjs[0]) * height/width
adjs[5] = a5    # horz pos of elbow wrt shape width, 0 is margin line

# rotate 45 degrees counter-clockwise
callout_sp.rotation = -45.0