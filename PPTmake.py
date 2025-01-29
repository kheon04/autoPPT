from pptx import Presentation
import shutil

# 파일 복사, 새로운 파일을 생성
destination_file_name = r"test1" # file_name
source = "Theme.pptx"  # Theme pptx path
destination_folder = r"lyrics\\" # destination folder
destination = destination_folder + destination_file_name + r".pptx"
shutil.copyfile(source, destination)

# pptx 파일 열기
prs = Presentation(destination)

# 새로운 슬라이드 만들기
# LAYOUT = 0번이 FADE, 1번이 NON_FADE
SLIDE_LAYOUT = prs.slide_layouts[0]
slide = prs.slides.add_slide(SLIDE_LAYOUT)

# Title에 제목을 저장
slide.shapes.title.text = "Born Again"
# SubTitle에 가사를 저장
tf = slide.placeholders[1].text_frame
tf.text = "주 안에서 내 영혼 다시 태어나" # 가사 입력
p = tf.add_paragraph() # Enter 입력
p.text = "이전 것은 지나고 새롭게 됐네" # 가사 입력
p = tf.add_paragraph() # Enter 입력
p = tf.add_paragraph() # Enter 입력

# pptx 파일 저장
prs.save(destination)