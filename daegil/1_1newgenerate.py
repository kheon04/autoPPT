from pptx import Presentation
import copy
import six
prs = Presentation()

layout = prs.slide_layouts[0]


def duplicate_slide(prs, index):
    template = prs.slides[index]
    try:
        blank_slide_layout = prs.slide_layouts[12]
    except:
        blank_slide_layout = prs.slide_layouts[len(prs.slide_layouts)]

    copied_slide = prs.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(
                value.reltype,
                value._target,
                value.rId
            )

    return copied_slide


path = "test.pptx"
prs_object = Presentation()
#첫번째 슬라이드를 소스 슬라이드로 가져옴
source_slide = prs.slides[0]
#새로 추가할 슬라이드의 레이아웃을 빈페이지[6]으로 정의
slide_layout = prs.slide_layouts[6]



















# def clone_slide(prs, slide):
#     new_slide = prs.slides.add_slide(slide.slide_layout)
    
#     # 원본 슬라이드의 모든 내용을 복사
#     for shape in slide.shapes:
#         if shape.has_text_frame:
#             new_shape = new_slide.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)
#             new_shape.text = shape.text

#     return new_slide

# def main():
#     # 프레젠테이션 객체 생성
#     prs = Presentation('your_prsentation.pptx')

#     # 복제할 슬라이드 인덱스 선택 (예: 0번째 슬라이드)
#     original_slide_index = 0
#     original_slide = prs.slides[original_slide_index]

#     # 선택한 슬라이드를 복제
#     new_slide = clone_slide(1, original_slide)

#     # 새로운 프레젠테이션 저장
#     prs.save('new_prsentation.pptx')

# if __name__ == "__main__":
#     main()