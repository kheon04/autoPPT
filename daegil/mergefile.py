from pptx import Presentation

# filepath = input("원본파일의 경로를 입력하십시오.")
# filepath = "C:\\.dev\\python\\test.pptx"


# f = open("C:\\.dev\\python\\test.pptx", "r")
# f.close()
# prs = Presentation("C:\\.dev\\python\\test.pptx")
def launchfile(path):
    try:
        prs = Presentation(path)

        print("opened.")
    except:
        prs = Presentation()
        prs.save(path)
        print("Success.")


def select_shape_by_text(slide, text):
    for x in slide.shapes:
    # 선택한 슬라이드 내 모든 shape들 중에서
        if x.has_text_frame and x.text == text:
        # text_frame이 있고 해당 shape의 text값이 입력한 text값과 같으면
            return x
            # 해당 shape을 반환합니다.
    print('요청한 Shape를 찾을 수 없습니다.')


import copy           
def copy_slide(prs, index):
    template = prs.slides[index]
    # prs 내 slide 중 index번째 슬라이드를 선택하여 template라는 변수명으로 저장합니다.
    try:
        blank_slide_layout = prs.slide_layouts.get_by_name('빈 화면')
    except:
        blank_slide_layout = prs.slide_layouts[0]
    # 새 슬라이드를 만들때 활용할 layout을 선택합니다.
    # '빈 화면'이라는 레이아웃을 선택하되 에러가 날 경우 layout중 첫번째를 선택합니다.
    
    copied_slide = prs.slides.add_slide(blank_slide_layout)
    # '빈 화면' 레이아웃을 적용한 새 슬라이드를 만들어 copied_slide라는 변수명으로
    # 저장합니다.
    
    for shape in template.shapes:
        elem = shape.element
        new_elem = copy.deepcopy(elem)
        copied_slide.shapes._spTree.insert_element_before(new_elem, 'p:extLst')
    # template 슬라이드에서 모든 shapes의 element를 복사하여
    # copied_slide에 붙여넣습니다.
        
    return copied_slide

filepath = "test.pptx"
launchfile(filepath)
