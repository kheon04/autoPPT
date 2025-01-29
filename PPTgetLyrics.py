from pptx import Presentation
import sqlite3


initDone = True
path = r"lyrics_ppt//"+"Born Again.pptx"  # Theme pptx path
prs = Presentation(path)

# Definition

def get_record_id(song_id, record_num):
    """
    ID를 생성합니다.
    f"{song_id:09}_{record_num:03}"
    """

    # ID 생성: songID_temp
    record_id = f"{song_id:09}_{record_num:03}"
    
    return record_id

def get_song_datas(slides, song_id):
    """
    song_Data 형태의 list를 return (?,?,?,?,?)
    prs.slides와 song_id를 인자로 받음
    """
    # slideNum, record_num = 0 초기화
    slide_num = 0
    record_num = 0

    # song_id 갱신
    song_id += 1
    for slide in slides:
        subTitle = slide.placeholders[1]
        # lines에 text를 '\n'을 기준으로 나누어 리스트형태로 저장
        lines = subTitle.text_frame.text
        lines = lines.split('\n')
        # 뒤에 존재하는 공백문자를 제거 (공백을 제거할지 기능을 설정할 것.)
        while (lines[len(lines)-1]==""):
            lines.pop()

        # 슬라이드의 가사가 2line일때
        if (len(lines)==2):
            # slide_num 갱신
            slide_num += 1

            # SongData의 초기 Signs 값은 None으로 저장함
            for line in lines:
                # record_num, record_id 갱신
                record_num += 1
                record_id = get_record_id(song_id=song_id, record_num=record_num)

                # songData 생성, append
                songData = (record_id, song_id, slide_num, None, line)
                songDataList.append(songData)
    return songDataList





# initDatas
try:
    # 데이터베이스에 연결
    conn = sqlite3.connect(r"databases\\songData.db")

    # 커서 생성
    cur = conn.cursor()
    cur.execute("""SELECT LastSongID FROM INIT""")
    LASTSONGID = cur.fetchone()[0]
    LASTSONGID = int(LASTSONGID)

    # 현재 문서에서 사용할 song_id 변수 생성
    song_id = LASTSONGID

finally:
    # 연결 종료 (예외가 발생해도 실행됨)
    conn.close()
    


# songData라는 입력할 데이터(?,?,?,?,?)의 list형태로 songDataList를 생성
songDataList = []
songDataList = get_song_datas(prs.slides, song_id)
print(tuple(songDataList))
LASTSONGID += 1


# addDatas
try:
    # 데이터베이스에 연결
    conn = sqlite3.connect(r"databases\\songData.db")

    # 커서 생성
    cur = conn.cursor()

    # 가사 Data 저장
    insert_sql = "INSERT INTO Main (RecordID, SongID, SlideNum, Signs, Line) VALUES (?, ?, ?, ?, ?);"
    cur.executemany(insert_sql, tuple(songDataList))

    # SongIDData Update
    song_id = LASTSONGID
    song_name = input("song_name을 입력하십시오: ")
    insert_sql = "INSERT INTO SongID (SongID, Name, [Modified Date], [Made Date]) VALUES (?, ?, datetime('now'), datetime('now'));"
    cur.execute(insert_sql, (LASTSONGID, song_name))

    # LastSongID update
    
    print(f"Song_ID update: updated to {LASTSONGID}!")
    cur.execute("UPDATE INIT SET LastSongID = ?", (int(LASTSONGID),))
    conn.commit()

finally:
    # 연결 종료 (예외가 발생해도 실행됨)
    conn.close()





class CCM:
    V, PC, C, B = "Verb" , "PreChorus", "Chorus", "Bridge"
    V1, V2, V3, V4, V5, V6 = "Verb1", "Verb2", "Verb3", "Verb4" ,"Verb5","Verb6"


